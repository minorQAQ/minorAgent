"""Agent 通用工具：会话历史、上下文压缩、文档解析与多模态辅助。

系统定位:
    - ``core/runtime``、``core/llm`` 依赖 summarizer、assistant_text、Documents_process
    - 所有历史数据基于 JSON 文件存储（multi_inputs 目录下）

主要子系统:
    1. 历史与元数据: get_history / get_session_meta / update_session_meta
    2. 上下文压缩: summarizer + 轻量/全量压缩
    3. 文档处理: Documents_process 多格式解析
    4. 文本/音频: split_text、merge_wavs、assistant_text

可扩展性:
    - 压缩策略阈值见 env_utils；可插拔向量记忆替代纯摘要
    - DOC_EXT_MAP 可扩展新文档类型处理器
"""
import json
import os
import shutil
from typing import List, Tuple, Any
import pathlib
import mimetypes, re

from langchain_core.chat_history import BaseChatMessageHistory
from langchain_core.messages import BaseMessage, SystemMessage, HumanMessage, AIMessage
from langchain_core.messages import messages_from_dict, messages_to_dict
from langchain_core.prompts import ChatPromptTemplate

from agent.memory.system_prompt import SUMMARY_SYSTEM_PROMPT
from agent.utils.image_utils import IMAGE_FILE_EXTENSIONS

# 音频文件扩展名（小写，含点）
AUDIO_FILE_EXTENSIONS = frozenset({".wav", ".mp3", ".flac", ".ogg", ".m4a", ".aac", ".wma", ".opus", ".aiff", ".ape", ".amr", ".au"})

# 扩展名到内部类型标签的映射
# plain_text 覆盖所有 UTF-8 文本/代码文件（与前端附件分类保持一致）
DOC_EXT_MAP = {
    ".txt": "plain_text", ".json": "plain_text", ".md": "plain_text", ".markdown": "plain_text",
    ".js": "plain_text", ".mjs": "plain_text", ".cjs": "plain_text", ".ts": "plain_text",
    ".tsx": "plain_text", ".jsx": "plain_text", ".vue": "plain_text", ".svelte": "plain_text",
    ".py": "plain_text", ".pyw": "plain_text", ".html": "plain_text", ".htm": "plain_text",
    ".css": "plain_text", ".scss": "plain_text", ".less": "plain_text", ".sass": "plain_text",
    ".xml": "plain_text", ".yaml": "plain_text", ".yml": "plain_text", ".toml": "plain_text",
    ".ini": "plain_text", ".cfg": "plain_text", ".conf": "plain_text", ".log": "plain_text",
    ".env": "plain_text", ".sh": "plain_text", ".bat": "plain_text", ".cmd": "plain_text",
    ".ps1": "plain_text", ".sql": "plain_text", ".r": "plain_text", ".go": "plain_text",
    ".rs": "plain_text", ".rb": "plain_text", ".php": "plain_text", ".swift": "plain_text",
    ".kt": "plain_text", ".scala": "plain_text", ".cpp": "plain_text", ".cc": "plain_text",
    ".c": "plain_text", ".h": "plain_text", ".hpp": "plain_text", ".java": "plain_text",
    ".m": "plain_text", ".svg": "plain_text",
    ".docx": "docx",
    ".docm": "docx",
    ".pptx": "pptx",
    ".pdf": "pdf",
    ".csv": "tabular",
    ".xlsx": "tabular",
    ".xls": "tabular",
}

# ---------- 目录与路径 ----------
_AGENT_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
HISTORY_ROOT = os.path.join(_AGENT_DIR, "history")
SESSIONS_ROOT = os.path.join(HISTORY_ROOT, "sessions")

# 确保 sessions 目录存在
os.makedirs(SESSIONS_ROOT, exist_ok=True)

# ---------- JSON 文件聊天历史实现 ----------
class JsonChatMessageHistory(BaseChatMessageHistory):
    """基于 turn 文件 的聊天历史，替代 SQLChatMessageHistory。

    系统定位:
        替代 SQLite 存储，直接从 sessions/{session_id}/turn_*.json
        读写 LangChain 消息。

    可扩展性:
        可增加消息索引、时间范围查询等。
    """

    def __init__(self, session_id: str):
        self.session_id = session_id
        self._messages: list[BaseMessage] = []
        self._loaded = False

    @property
    def messages(self) -> list[BaseMessage]:
        if not self._loaded:
            self._load()
        return list(self._messages)

    def _load(self) -> None:
        """从 turn 文件加载消息。"""
        self._loaded = True
        self._messages.clear()

        turn_msgs = load_all_turn_messages(self.session_id)
        for item in turn_msgs:
            msg = self._dict_to_langchain(item)
            if msg is not None:
                self._messages.append(msg)

    def add_message(self, message: BaseMessage) -> None:
        """追加一条消息到内存（消息由 turn 级存储管理）。"""
        if not self._loaded:
            self._load()
        self._messages.append(message)

    @staticmethod
    def _dict_to_langchain(item: dict) -> BaseMessage | None:
        """将消息字典转为 LangChain BaseMessage。"""
        if not isinstance(item, dict):
            return None
        role = item.get("role")
        if role == "user":
            user_data = item.get("user") or {}
            content = user_data.get("content")
            if content is None:
                return None
            text_parts = []
            if isinstance(content, list):
                for p in content:
                    if isinstance(p, dict) and p.get("type") == "text":
                        text_parts.append(p.get("text") or "")
                text = "\n".join(text_parts) if text_parts else ""
            else:
                text = str(content) if content else ""
            if not text:
                return None
            return HumanMessage(content=text)
        elif role == "assistant":
            asst_data = item.get("assistant") or {}
            text = (asst_data.get("text") or "").strip()
            return AIMessage(content=text) if text else None
        return None

    def add_messages(self, messages: list[BaseMessage]) -> None:
        """批量追加消息。"""
        for msg in messages:
            self.add_message(msg)

    def clear(self) -> None:
        """清空内存消息（持久化数据由 turn 级存储管理，详见 delete_turns_after）。"""
        self._messages.clear()


# ---------- 历史记录操作 ----------
def get_history(session_id: str) -> JsonChatMessageHistory:
    """获取指定会话的 JSON 聊天历史对象。

    输入:
        session_id: 会话唯一标识。

    输出:
        JsonChatMessageHistory 实例，底层为 turn 文件。

    系统定位: summarizer、save_messages_to_history 等读写对话记录。
    """
    return JsonChatMessageHistory(session_id=session_id)


# ---------- session_meta 元数据管理 ----------
def _load_session_meta(session_id: str) -> dict:
    """加载整个 session_meta 数据（后端无关：mysql 存 sessions.extra.agent_meta，json 存文件）。"""
    from agent.history import session_storage
    return session_storage.load_session_extra(session_id).get("agent_meta", {})


def _save_session_meta(session_id: str, meta: dict) -> None:
    """保存整个 session_meta 数据（后端无关，不覆盖 extra 中的 tokens）。"""
    from agent.history import session_storage
    existing = session_storage.load_session_extra(session_id) or {}
    existing["agent_meta"] = meta
    session_storage.save_session_extra(session_id, existing)


def get_session_meta(session_id: str, agent_name: str) -> dict:
    """读取会话在指定 Agent 下的压缩元数据（不含 turn_tool_context）。"""
    meta = _load_session_meta(session_id)
    agent_meta = meta.get(agent_name, {})
    return {
        "cursor_index": agent_meta.get("cursor_index", -1),
        "compressed_content": agent_meta.get("compressed_content", ""),
    }


def update_session_meta(session_id: str, agent_name: str,
                        cursor_index: int, compressed_content: str) -> None:
    """写入或更新 session_meta（仅 cursor 和压缩内容）。"""
    meta = _load_session_meta(session_id)
    meta[agent_name] = {
        "cursor_index": cursor_index,
        "compressed_content": compressed_content,
    }
    _save_session_meta(session_id, meta)


# ---------- 会话级思考档位（存于 session_meta.json 顶层） ----------
def get_session_thinking_level(session_id: str) -> str:
    """读取会话级思考档位（顶层 thinking_level 键）。

    输入:
        session_id: 会话 ID。

    输出:
        档位字符串（low | high | xhigh | max | ultra）；未配置返回 ""。
    """
    if not session_id:
        return ""
    meta = _load_session_meta(session_id)
    return str(meta.get("thinking_level", "") or "")


def set_session_thinking_level(session_id: str, level: str) -> None:
    """写入会话级思考档位（顶层 thinking_level 键，不覆盖 agent_meta / tokens）。"""
    if not session_id:
        return
    meta = _load_session_meta(session_id)
    meta["thinking_level"] = level
    _save_session_meta(session_id, meta)


# ---------- 动态尾部长期记忆（TodoList 状态 + 循环提醒） ----------
def update_session_dynamic_tail(session_id: str, agent_key: str, text: str, max_entries: int = 100) -> None:
    """持久化当前轮的动态尾部提示（TodoList 状态 + 循环提醒）为长期记忆。

    输入:
        session_id: 会话 ID。
        agent_key: session_meta 中的 agent 键（主 Agent 为 "main"，子 Agent 用其名称）。
        text: 本轮动态尾部文本。
        max_entries: 历史条目上限（防止无界增长），超出丢弃最旧。

    说明:
        与 compressed_content 合并保存（不覆盖）；连续相同内容去重，
        使 todo 进度等状态变更形成可追溯的历史记录，供后续轮次注入查看
        （充分利用前缀缓存：变化部分始终位于消息末尾）。
    """
    if not session_id or not agent_key or not text:
        return
    meta = _load_session_meta(session_id)
    entry = dict(meta.get(agent_key) or {})
    history = list(entry.get("dynamic_tail_history") or [])
    if history and history[-1] == text:
        return  # 去重：内容未变化则不重复记录
    history.append(text)
    if len(history) > max_entries:
        history = history[-max_entries:]
    entry["dynamic_tail_history"] = history
    meta[agent_key] = entry
    _save_session_meta(session_id, meta)


def get_session_dynamic_tail_history(session_id: str, agent_key: str) -> list[str]:
    """读取指定 Agent 的动态尾部长期记忆历史（时间正序，最新在末尾）。

    输出:
        [tail_text, ...]；无记录返回空列表。
    """
    if not session_id or not agent_key:
        return []
    meta = _load_session_meta(session_id)
    entry = meta.get(agent_key) or {}
    return list(entry.get("dynamic_tail_history") or [])


def build_tool_history_messages(session_id: str, agent_key: str = "main") -> list:
    """构建压缩游标之后的"历史思考与工具调用"合成消息列表。

    输入:
        session_id: 会话 ID。
        agent_key: session_meta 中的 agent 键（主 Agent 为 "main"）。

    输出:
        synthetic HumanMessage（带 ``tool_ctx`` 标记）列表；无历史或
        游标之后无内容时为空列表。

    上下文拼接格式（前缀缓存友好，确定性构建）:
        [历史思考] {thinking}     ← 该工具调用前模型产出的深度思考（reasoning_content）
        [历史反思] {reflection}   ← 该工具调用前模型产出的反思（content）
        [历史工具调用] {name}({args}) → {result}
        （思考、反思、工具调用按因果顺序交错；思考/反思独立注入——工具无
          结果文本（失败/无返回值）时仍保留；每轮末尾未被工具调用消费的
          尾部思考/反思去重后追加在该轮工具调用之后）

    系统定位:
        runtime.execute_agent 构建上下文时调用，替代原"最近5轮"注入，
        改为注入压缩游标之后的所有轮次思考+工具调用历史（游标之前的
        部分已包含在压缩摘要中，由 ``compressed_content`` 承载）。
        所有消息从磁盘确定性构建、id 稳定（think_ctx_{idx}/
        reflect_ctx_{idx}/tool_ctx_{idx}），满足 LLM 前缀缓存条件；
        compress 节点按 id 精确移除。
    """
    if not session_id:
        return []
    meta = get_session_meta(session_id, agent_key)
    cursor = meta.get("cursor_index", -1)
    result: list = []
    try:
        from agent.history.tool_call_recorder import list_turns
        turns = list(list_turns(session_id, limit=100000))  # 时间倒序
        turns.reverse()  # 转为时间正序
        idx = 0
        for turn in turns:
            # 本轮已被工具调用消费的思考/反思内容（用于尾部去重；
            # 含游标之前的工具调用——被摘要覆盖的思考同样不得重复注入）
            consumed_thinkings: set[str] = set()
            consumed_reflections: set[str] = set()
            for tc in turn.get("tool_calls") or []:
                _think = str(tc.get("thinking", "") or "").strip()
                _reflect = str(tc.get("reflection", "") or "").strip()
                if _think:
                    consumed_thinkings.add(_think)
                if _reflect:
                    consumed_reflections.add(_reflect)
                if idx >= cursor:
                    tn = tc.get("name", "unknown")
                    ta = json.dumps(tc.get("args", {}), ensure_ascii=False)[:300]
                    tr = str(tc.get("result_text", "") or "")[:500]
                    # 思考/反思独立注入：即使工具无结果文本（失败/无返回值），
                    # 模型的思考/反思仍有上下文价值
                    if _think:
                        result.append(HumanMessage(
                            content=f"[历史思考] {_think[:500]}",
                            additional_kwargs={"synthetic": True, "tool_ctx": True},
                            id=f"think_ctx_{idx}",
                        ))
                    if _reflect:
                        result.append(HumanMessage(
                            content=f"[历史反思] {_reflect[:500]}",
                            additional_kwargs={"synthetic": True, "tool_ctx": True},
                            id=f"reflect_ctx_{idx}",
                        ))
                    if tn and tr:
                        # 分配稳定 id（基于全局工具调用序号），供压缩节点 RemoveMessage 精确移除
                        result.append(HumanMessage(
                            content=f"[历史工具调用] {tn}({ta}) → {tr}",
                            additional_kwargs={"synthetic": True, "tool_ctx": True},
                            id=f"tool_ctx_{idx}",
                        ))
                idx += 1
            # 尾部思考/反思：该轮末尾未被任何工具调用消费的（分别去重后追加）
            if idx > cursor:
                think_tail_idx = 0
                reflect_tail_idx = 0
                for ref in turn.get("reflections") or []:
                    _thought = str(ref.get("thought", "") or "").strip()
                    _reflection = str(ref.get("reflection", "") or "").strip()
                    if _thought and _thought not in consumed_thinkings:
                        result.append(HumanMessage(
                            content=f"[历史思考] {_thought[:500]}",
                            additional_kwargs={"synthetic": True, "tool_ctx": True},
                            id=f"think_tail_{turn.get('turn_id', idx)}_{think_tail_idx}",
                        ))
                        think_tail_idx += 1
                    if _reflection and _reflection not in consumed_reflections:
                        result.append(HumanMessage(
                            content=f"[历史反思] {_reflection[:500]}",
                            additional_kwargs={"synthetic": True, "tool_ctx": True},
                            id=f"reflect_tail_{turn.get('turn_id', idx)}_{reflect_tail_idx}",
                        ))
                        reflect_tail_idx += 1
                    # 旧格式兼容（仅 content，无 thought/reflection 字段）：
                    # 视为思考注入（与旧行为一致）
                    if not _thought and not _reflection:
                        rc = str(ref.get("content", "") or "").strip()
                        if rc and rc not in consumed_thinkings:
                            result.append(HumanMessage(
                                content=f"[历史思考] {rc[:500]}",
                                additional_kwargs={"synthetic": True, "tool_ctx": True},
                                id=f"think_tail_{turn.get('turn_id', idx)}_{think_tail_idx}",
                            ))
                            think_tail_idx += 1
    except Exception:
        pass
    return result


# ---------- Turn 级消息存储 ----------
import threading as _threading

# Turn JSON 文件前缀/扩展名
TURN_JSON_PREFIX = "turn_"


def turn_files_dir(session_id: str, turn_id: str) -> str:
    """返回指定 turn 的附件目录路径。"""
    return os.path.join(SESSIONS_ROOT, session_id, turn_id, "files")


def save_turn_messages(session_id: str, turn_id: str, messages: list[dict]) -> None:
    """将一轮对话的消息持久化（后端无关：mysql 只入库，json 只写文件，不双写）。

    系统定位:
        消息持久化的唯一入口（sync_turn_from_chat_history 调用）。
        附件二进制文件始终落盘（由调用方处理），本函数仅负责对话消息记录的持久化。
    """
    from agent.history import session_storage
    session_storage.save_turn(session_id, turn_id, messages)


def load_all_turn_messages(session_id: str) -> list[dict]:
    """按时间顺序加载所有 turn 的消息，返回扁平列表（后端无关）。"""
    from agent.history import session_storage
    return session_storage.load_messages(session_id)


def delete_turns_after(session_id: str, keep_until_turn_id: str) -> int:
    """删除指定 turn_id 之后（按文件名排序）的所有 turn 记录并清理对应 artifact 目录。

    记录删除（消息 + 工具调用）由统一存储处理后端无关；artifact 目录（以 turn_id
    命名的目录，两个后端均落盘）在此清理。返回删除的 turn 数。

    注意：keep_until_turn_id 自身的 turn 记录会被保留（只删除严格大于它的）。
    """
    sess_dir = os.path.join(SESSIONS_ROOT, session_id)
    if os.path.isdir(sess_dir):
        for dname in sorted(os.listdir(sess_dir)):
            dp = os.path.join(sess_dir, dname)
            if os.path.isdir(dp) and dname > keep_until_turn_id:
                shutil.rmtree(dp, ignore_errors=True)
    from agent.history import session_storage
    return session_storage.delete_turns_after(session_id, keep_until_turn_id)


def get_last_turn_id(session_id: str) -> str | None:
    """获取最后一个 turn_id（后端无关）；无记录返回 None。"""
    from agent.history import session_storage
    return session_storage.get_last_turn_id(session_id) or None


def session_dir(session_id: str) -> str:
    """返回会话目录路径（所有 turn 和附件根目录）。"""
    return os.path.join(SESSIONS_ROOT, session_id)


# 跨请求共享 pending turn_id
_pending_turn_ids: dict[str, str] = {}
_pending_lock = _threading.Lock()


def set_pending_turn_id(session_id: str, turn_id: str) -> None:
    with _pending_lock:
        _pending_turn_ids[session_id] = turn_id


def take_pending_turn_id(session_id: str) -> str | None:
    """取出并清除 pending turn_id（仅取一次）。"""
    with _pending_lock:
        return _pending_turn_ids.pop(session_id, None)


# ---------- 共享压缩摘要函数 ----------
def summarize_chat_text(llm, chat_text: str) -> str:
    """用 LLM 对对话文本进行压缩摘要，供 in-loop 压缩和手动压缩复用。

    输入:
        llm: 压缩用 LLM 实例。
        chat_text: 待压缩的对话文本。

    输出:
        摘要字符串。

    系统定位:
        nodes.maybe_compress_context（图内自动压缩节点）的摘要生成核心。
    """
    prompt = ChatPromptTemplate.from_messages([
        ("system", SUMMARY_SYSTEM_PROMPT),
        ("human", "请对以下对话片段进行压缩摘要，保留任务目标、已完成的工作、关键发现和当前待办：\n\n{chat}"),
    ])
    chain = prompt | llm
    result = chain.invoke({"chat": chat_text}, config={"metadata": {"output_type": "text"}})
    return result.content


# ---------- 辅助函数 ----------
def assistant_text(value) -> str:
    """将 LLM 响应的 content/reasoning 等多形态字段归一为可见纯文本。

    输入:
        value: str | list | dict | None 等网关返回结构。

    输出:
        去首尾空白后的字符串；无法解析时 str(value) 或 JSON。

    系统定位:
        runtime、llm、nodes 统一提取助手回复与 reasoning。

    可扩展性:
        新增厂商字段时在 key 列表中补充即可。
    """
    if value is None:
        return ""
    if isinstance(value, str):
        return value.strip()
    if isinstance(value, list):
        parts: list[str] = []
        for item in value:
            if isinstance(item, str) and item.strip():
                parts.append(item.strip())
            elif isinstance(item, dict):
                for key in ("text", "reasoning_content", "content"):
                    v = item.get(key)
                    if isinstance(v, str) and v.strip():
                        parts.append(v.strip())
                        break
        return "\n".join(parts).strip()
    if isinstance(value, dict):
        for key in ("text", "content", "input", "output", "query", "answer", "result"):
            v = value.get(key)
            if isinstance(v, str) and v.strip():
                return v.strip()
        try:
            return json.dumps(value, ensure_ascii=False)
        except TypeError:
            return str(value).strip()
    return str(value).strip()

# ---------- 附件描述文本 ----------
def make_attachment_text(category: str, file_path: str, extra: str = "") -> str:
    """生成用户提供的附件描述文本，供 LLM 上下文注入。

    输入:
        category: 附件类别，如 "音频"、"图片"、"文件"、"附件文件"。
        file_path: 文件路径（已规范化）。
        extra: 附加信息，如尺寸 "1920×1080px"、ASR 文本等。

    输出:
        格式化后的提示文本字符串。
    """
    normalized = file_path.replace("\\", "/")
    if category == "图片" and extra:
        return f"[用户提供了图片，文件路径: {normalized}，尺寸: {extra}]"
    if category == "图片":
        return f"[用户提供了图片，文件路径: {normalized}]"
    if category == "音频" and extra:
        return f"[用户提供了音频，文件路径: {normalized}]\n{extra}"
    if extra:
        return extra  # 已有完整格式（如 ASR 失败、文档内容等）
    if category == "附件文件":
        return f"[用户提供了附件文件: {normalized}]"
    return f"[用户提供了{category}，文件路径: {normalized}]"


class Documents_process:
    """多格式用户附件的分类、解析与多模态内容提取。

    系统定位:
        core/llm._preprocess_messages 处理用户上传文件时调用 process()。

    可扩展性:
        新增格式：扩展 DOC_EXT_MAP 并实现 process_xxx 方法。
    """

    @staticmethod
    def doc_classify(file_path: str) -> str:
        """根据扩展名与 MIME 返回内部类型标签。

        输入:
            file_path: 本地文件路径。

        输出:
            plain_text | docx | pdf | image | tabular。

        系统定位:
            process() 分发入口。

        可扩展性:
            不支持类型抛出 ValueError，便于上层捕获并提示用户。
        """
        ext = pathlib.Path(file_path).suffix.lower()
        # 先检查文档映射表
        if ext in DOC_EXT_MAP:
            return DOC_EXT_MAP[ext]

        # 再检查是否为图片
        if ext in IMAGE_FILE_EXTENSIONS:
            return "image"

        # 再检查是否为音频
        if ext in AUDIO_FILE_EXTENSIONS:
            return "audio"
        
        # 最后使用 MIME 类型判断
        mime, _ = mimetypes.guess_type(file_path)
        if mime:
            if mime.startswith("text/"):
                return "plain_text"
            if mime.startswith("image/"):
                return "image"
            if mime.startswith("audio/"):
                return "audio"

        raise ValueError(f"不支持的文件类型: {file_path}")

    # ---------- 各类文件的具体处理方法 ----------
    @staticmethod
    def process_plain_text(file_path: str) -> tuple[str, list]:
        """读取 UTF-8 纯文本文件。

        输入: file_path 本地路径。
        输出: (文本内容, 空图片列表)。
        系统定位: plain_text 类型处理。
        可扩展性: 可增加编码探测（chardet）。
        """
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read(), []

    @staticmethod
    def process_docx(file_path: str) -> tuple[str, list]:
        """按文档顺序提取 Word 文本、内嵌图片与表格。

        输入: file_path .docx 路径。
        输出: (带 <image>N</image>、<table> 标记的文本, 图片字节列表)。
        系统定位: 用户上传 docx 附件解析。
        可扩展性: 依赖 python-docx；可扩展页眉页脚、脚注。
        """
        try:
            from docx import Document
            from docx.oxml.ns import qn
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
        except ImportError:
            raise ImportError("需要安装 python-docx: pip install python-docx")

        doc = Document(file_path)
        content_parts = []
        image_bytes_list = []
        image_counter = 0

        # ---------- 辅助函数：按文档顺序迭代段落和表格 ----------
        def iter_block_items(parent):
            """生成文档中的段落和表格，保持原始顺序"""
            from docx.oxml.ns import qn as ns_qn
            if hasattr(parent, 'element'):
                parent_elm = parent.element.body
            else:
                parent_elm = parent

            # 建立元素到段落的映射，便于 O(1) 查找
            para_map = {p._element: p for p in doc.paragraphs}
            # 表格同样建立映射
            table_map = {t._element: t for t in doc.tables}

            for child in parent_elm.iterchildren():
                if child.tag == ns_qn('w:p'):
                    para = para_map.get(child)
                    if para is not None:
                        yield para
                elif child.tag == ns_qn('w:tbl'):
                    table = table_map.get(child)
                    if table is not None:
                        yield table

        def extract_image_from_run(run, para):
            """从 run 中提取图片，返回 <image>N</image> 字符串和图片二进制列表"""
            nonlocal image_counter
            drawings = run._element.findall('.//' + qn('w:drawing'))
            imgs = []
            for drawing in drawings:
                blips = drawing.findall('.//' + qn('a:blip'))
                for blip in blips:
                    embed = blip.get(qn('r:embed'))
                    if embed:
                        image_counter += 1
                        related_part = para.part.related_parts[embed]
                        image_bytes_list.append(related_part.blob)
                        imgs.append(f"<image>{image_counter}</image>")
            return imgs

        pending_title = ""   # 缓存可能存在的图表短标题

        for item in iter_block_items(doc):
            if hasattr(item, 'text'):   # 段落
                para = item
                para_parts = []
                # 提取图片
                for run in para.runs:
                    para_parts.extend(extract_image_from_run(run, para))
                # 提取文字
                text = para.text.strip()
                if text:
                    para_parts.insert(0, text)
                para_str = "".join(para_parts)

                if not para_str:
                    continue

                # 判断是否是短标题（如“表1 一个表格”、“图1 一张图片”）
                if len(text) < 30 and (text.startswith("表") or text.startswith("图") or text.startswith("Table") or text.startswith("Fig")):
                    # 暂存为候选标题，等待后续表格/图片
                    pending_title = para_str
                    continue

                # 普通段落：先输出可能遗留的标题（没有后续图表匹配的），再输出本段
                if pending_title:
                    content_parts.append(pending_title + "\n")
                    pending_title = ""
                content_parts.append(para_str + "\n")

            else:   # 表格
                table = item
                table_str_parts = [f"===== Table ??? =====\n"]  # 标题稍后合并
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_str_parts.append("| " + " | ".join(row_data) + " |\n")
                table_content = "".join(table_str_parts)

                # 如果有暂存标题，合并到一起
                if pending_title:
                    # 替换表格占位标题
                    final_table = f"<table>{pending_title}\n{table_content[table_content.index('\n')+1:]}</table>"
                    content_parts.append(final_table + "\n")
                    pending_title = ""
                else:
                    # 没有标题，保留原始表格内容
                    content_parts.append(f"<table>{table_content}</table>\n")

        # 文件末尾的遗留标题
        if pending_title:
            content_parts.append(pending_title + "\n")

        return "".join(content_parts), image_bytes_list

    @staticmethod
    def process_pdf(file_path: str) -> tuple[str, list]:
        """提取 PDF 文本与内嵌图片占位符。

        输入: file_path .pdf 路径。
        输出: (分页拼接文本含 <image>N, 图片字节列表)。
        系统定位: PDF 附件入模前解析。
        可扩展性: 依赖 pypdf；扫描件可接 OCR 管线。
        """
        try:
            from pypdf import PdfReader 
        except ImportError:
            raise ImportError("需要安装 pypdf:pip install pypdf")
        
        reader = PdfReader(file_path)
        text_parts = []
        image_bytes_list = []
        image_counter = 0  # 全局图片计数器

        for page in reader.pages:
            page_content = []
            # 提取文字
            page_text = page.extract_text()
            if page_text:
                page_content.append(page_text)

            # 提取图片，并在文本中追加占位符
            for _ in page.images:
                image_counter += 1
                page_content.append(f"<image>{image_counter}</image>")

            # 收集图片二进制数据
            for image_file_object in page.images:
                image_bytes_list.append(image_file_object.data)

            # 合并当前页的内容
            if page_content:
                text_parts.append("".join(page_content))

        return "\n".join(text_parts), image_bytes_list

    @staticmethod
    def process_pptx(file_path: str) -> tuple[str, list]:
        """将 PPTX 转为文本 + SVG 代码 + 提取内嵌图片。

        使用 python-pptx 提取文字（可靠），同时用 ppt_utils 渲染 SVG。
        返回的文本中 base64 data URL 会被替换为占位符。

        输入: file_path .pptx 路径。
        输出: (按 slide 分隔的「文字 + SVG」文本, 图片字节列表)。
        """
        import base64
        import re
        from pathlib import Path
        from pptx import Presentation
        from agent.utils.ppt_utils import convert_pptx_to_svg

        pptx_path = Path(file_path)

        # ---- 1. python-pptx 提取文字（稳定可靠） ----
        prs = Presentation(str(pptx_path))
        text_parts: list[str] = []
        slide_texts: dict[int, list[str]] = {}  # slide_index -> text lines
        for si, slide in enumerate(prs.slides):
            lines: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = para.text.strip()
                        if txt:
                            lines.append(txt)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            lines.append(row_text)
            slide_texts[si] = lines
            if lines:
                text_parts.append(f"<slide>{si + 1} 文字></slide>\n" + "\n".join(lines))

        text_content = "\n\n".join(text_parts) if text_parts else "(未提取到文字内容)"

        # ---- 2. SVG 渲染 ----
        try:
            result = convert_pptx_to_svg(pptx_path)
        except Exception:
            # SVG 解析失败时仅返回文字
            return text_content, []

        # 收集所有图片字节
        image_bytes_list: list[bytes] = []
        seen = set()

        # 2a) 从 media_files 收集
        for art in result.slides:
            for blob in art.media_files.values():
                h = hash(blob)
                if h not in seen:
                    seen.add(h)
                    image_bytes_list.append(blob)

        # 2b) 从 SVG 中的 base64 data URL 提取图片字节，替换为占位符
        data_url_pat = re.compile(r'href="data:image/\w+;base64,([A-Za-z0-9+/=]+)"')
        img_idx = len(image_bytes_list)

        def _replace_data_url(m: re.Match) -> str:
            nonlocal img_idx
            try:
                blob = base64.b64decode(m.group(1))
                h = hash(blob)
                if h not in seen:
                    seen.add(h)
                    image_bytes_list.append(blob)
                    img_idx = len(image_bytes_list) - 1
                else:
                    for i, existing in enumerate(image_bytes_list):
                        if existing == blob:
                            img_idx = i
                            break
            except Exception:
                pass
            return f'href="#[内嵌图片 {img_idx + 1}]"'

        for art in result.slides:
            art.svg = data_url_pat.sub(_replace_data_url, art.svg)

        # ---- 3. 组装输出：文字 + SVG ----
        svg_parts: list[str] = []
        for art in result.slides:
            idx = art.index - 1  # 0-based
            section = f"<slide>{art.index}</slide>"
            if idx in slide_texts and slide_texts[idx]:
                section += "\n[文字]\n" + "\n".join(slide_texts[idx]) + "\n"
            section += "\n[布局]\n" + art.svg.strip()
            svg_parts.append(section)

        text = "\n\n".join(svg_parts)
        if not text.strip():
            text = text_content

        return text, image_bytes_list

    @staticmethod
    def process_tabular(file_path: str) -> tuple[str, list]:
        """将 CSV / xlsx 转为 <table> 标记的表格文本。

        输入: file_path .csv 或 .xlsx。
        输出: (多 sheet 表格字符串, 空图片列表)。
        系统定位: 结构化数据附件供 LLM 阅读。
        可扩展性: 可支持 xls、parquet；大表可采样行数。
        """
        try:
            import csv
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("需要安装 csv和openpyxl:pip install csv openpyxl")

        result_parts = []
        ext = pathlib.Path(file_path).suffix.lower()

        if ext == '.csv':
            # CSV 只有一个“sheet”
            with open(file_path, 'r', encoding='utf-8-sig', newline='') as f:
                reader = csv.reader(f)
                rows = list(reader)
            table_str = '\n'.join([' | '.join(row) for row in rows])
            result_parts.append(f"<table>Sheet1\n{table_str}</table>")

        elif ext == '.xlsx':
            wb = load_workbook(file_path, read_only=True, data_only=True)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append(' | '.join([str(cell) if cell is not None else '' for cell in row]))
                table_str = '\n'.join(rows)
                result_parts.append(f"<table>{sheet_name}\n{table_str}</table>")
            wb.close()

        elif ext == '.xls':
            import xlrd
            wb = xlrd.open_workbook(file_path)
            for sheet_name in wb.sheet_names():
                sheet = wb.sheet_by_name(sheet_name)
                rows = []
                for r in range(sheet.nrows):
                    row_vals = []
                    for c in range(sheet.ncols):
                        v = sheet.cell_value(r, c)
                        # 日期转可读字符串
                        if sheet.cell_type(r, c) == xlrd.XL_CELL_DATE:
                            try:
                                v = xlrd.xldate_as_datetime(v, wb.datemode)
                            except Exception:
                                pass
                        row_vals.append(str(v) if v is not None else '')
                    rows.append(' | '.join(row_vals))
                table_str = '\n'.join(rows)
                result_parts.append(f"<table>{sheet_name}\n{table_str}</table>")

        else:
            raise ValueError(f"Unsupported tabular file extension: {ext}")

        return "\n".join(result_parts), []

    @staticmethod
    def process_image(file_path: str) -> tuple[str, list]:
        """读取图片文件为字节，文本部分为空。

        输入: file_path 图片路径。
        输出: ("", [图片字节])。
        系统定位: 与 llm 中 image_url 分支配合多模态输入。
        可扩展性: 可在此预调用 resize_image_if_needed。
        """
        with open(file_path, "rb") as f:
            return "", [f.read()]

    @staticmethod
    def process_audio(file_path: str) -> tuple[str, list]:
        """读取音频文件为字节，文本部分为空（ASR 由调用方按需执行）。

        输入: file_path 音频路径。
        输出: ("", [音频字节])。
        系统定位: doc_tool / send_file 读取音频原始数据。
        可扩展性: 可在此执行音量标准化、降噪等预处理。
        """
        with open(file_path, "rb") as f:
            return "", [f.read()]

    # ---------- 统一处理入口 ----------
    @staticmethod
    def process(file_path: str) -> tuple[str, list]:
        """统一附件处理入口：分类后分发到对应 process_* 方法。

        输入: file_path 本地文件路径。
        输出: (文本内容, 图片字节列表)。
        系统定位: llm._preprocess_messages 文件 part 处理。
        可扩展性: 新增类型在 doc_classify 与分支中注册。
        """
        file_type = Documents_process.doc_classify(file_path)

        if file_type == "plain_text":
            return Documents_process.process_plain_text(file_path)
        elif file_type == "docx":
            return Documents_process.process_docx(file_path)
        elif file_type == "pptx":
            return Documents_process.process_pptx(file_path)
        elif file_type == "pdf":
            return Documents_process.process_pdf(file_path)
        elif file_type == "image":
            return Documents_process.process_image(file_path)
        elif file_type == "audio":
            return Documents_process.process_audio(file_path)
        elif file_type == "tabular":
            return Documents_process.process_tabular(file_path)
        else:
            raise ValueError(f"未处理的文件类型: {file_type}")

def audio_file_to_text(audio_path: str) -> str:
    """将音频文件通过 ASR 服务转为文本。

    输入:
        audio_path: 本地音频文件路径（支持 wav/mp3/flac/ogg/m4a 等）。

    输出:
        识别出的文本字符串。

    系统定位:
        doc_tool / send_file 等工具通用的音频转文字能力。

    可扩展性:
        可增加语言参数、说话人分离等高级 ASR 选项。
    """
    import os as _os
    import tempfile as _tempfile
    import base64 as _base64
    import requests as _requests

    from agent.utils import env_utils as _env_utils

    if not _os.path.exists(audio_path):
        raise FileNotFoundError(f"音频文件不存在: {audio_path}")

    ext = _os.path.splitext(audio_path)[1].lower()
    wav_path = audio_path
    cleanup_wav = False

    if ext != ".wav":
        try:
            from pydub import AudioSegment
            audio = AudioSegment.from_file(audio_path)
            fd, wav_path = _tempfile.mkstemp(suffix=".wav", prefix="asr_")
            _os.close(fd)
            audio.export(wav_path, format="wav")
            cleanup_wav = True
        except ImportError:
            raise RuntimeError("非 WAV 音频需要安装 pydub 和 ffmpeg: pip install pydub")

    asr_cfg = _env_utils.get_service_model_config("asr")
    try:
        with open(wav_path, "rb") as f:
            audio_bytes = f.read()
        b64 = _base64.b64encode(audio_bytes).decode("ascii")
        data_url = f"data:audio/wav;base64,{b64}"
        payload = {
            "model": asr_cfg["model"],
            "messages": [{
                "role": "user",
                "content": [{"type": "audio_url", "audio_url": {"url": data_url}}]
            }]
        }
        endpoint = f"{asr_cfg['base_url']}/v1/chat/completions"
        headers = {"Content-Type": "application/json"}
        if asr_cfg["api_key"]:
            headers["Authorization"] = f"Bearer {asr_cfg['api_key']}"
        resp = _requests.post(
            endpoint,
            headers=headers,
            json=payload,
            timeout=120,
        )
        resp.raise_for_status()
        result = resp.json()
        return result["choices"][0]["message"]["content"].strip()
    finally:
        if cleanup_wav and _os.path.exists(wav_path):
            _os.remove(wav_path)

    